import os
import pandas as pd
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from PIL import Image
import pytesseract


class UnifiedDataLoader:

    def load(self, file_path: str):
        ext = os.path.splitext(file_path)[1].lower()

        loaders = {
            ".pdf": self._load_pdf,
            ".doc": self._load_docx,
            ".docx": self._load_docx,
            ".ppt": self._load_ppt,
            ".pptx": self._load_ppt,
            ".xls": self._load_excel,
            ".xlsx": self._load_excel,
            ".jpg": self._load_image,
            ".jpeg": self._load_image,
            ".png": self._load_image,
        }

        if ext not in loaders:
            raise ValueError(f"Unsupported file type: {ext}")

        return loaders[ext](file_path)

    def _load_pdf(self, file_path):
        reader = PdfReader(file_path)
        return [
            {
                "text": page.extract_text() or "",
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "pdf",
                    "page": i + 1,
                },
            }
            for i, page in enumerate(reader.pages)
        ]

    def _load_docx(self, file_path):
        doc = Document(file_path)
        return [{
            "text": "\n".join(p.text for p in doc.paragraphs),
            "metadata": {
                "source": os.path.basename(file_path),
                "type": "docx",
            },
        }]

    def _load_ppt(self, file_path):
        prs = Presentation(file_path)
        docs = []

        for i, slide in enumerate(prs.slides):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text")
            )
            docs.append({
                "text": text,
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "pptx",
                    "slide": i + 1,
                },
            })

        return docs

    def _load_excel(self, file_path):
        sheets = pd.read_excel(file_path, sheet_name=None)
        return [
            {
                "text": df.to_string(index=False),
                "metadata": {
                    "source": os.path.basename(file_path),
                    "type": "excel",
                    "sheet": name,
                },
            }
            for name, df in sheets.items()
        ]

    def _load_image(self, file_path):
        image = Image.open(file_path)
        return [{
            "text": pytesseract.image_to_string(image),
            "metadata": {
                "source": os.path.basename(file_path),
                "type": "image",
            },
        }]
